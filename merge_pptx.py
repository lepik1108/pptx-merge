import glob, os
from spire.presentation import *
from pptx import Presentation as P

# Function to merge multiple presentations sequentially into a new presentation
def merge_presentations(presentation_files, output_file):
    # Create a new presentation
    merged_presentation = Presentation()
    # Remove the default empty slide in the new presentation
    merged_presentation.Slides.RemoveAt(0)

    for file in presentation_files:
        # Load each presentation
        presentation = Presentation()
        presentation.LoadFromFile(file)
        for slide in presentation.Slides:
            # Copy each slide to the new presentation
            merged_presentation.Slides.AppendBySlide(slide)

    # Save the new presentation
    merged_presentation.SaveToFile(output_file, FileFormat.Pptx2016)

# PowerPoint files_to_merge
#files_to_merge = ["p1.pptx", "p2.pptx", "p3.pptx"]

ppt_to_merge_folder="./presentations"
result_filename="merged.pptx"

os.chdir(ppt_to_merge_folder)

files_to_merge = sorted(
    [f for f in glob.glob("*.pptx") if f != result_filename]
)

print(f"Processing files {ppt_to_merge_folder}/:{files_to_merge}...")

# Call the function to merge the files sequentially
merge_presentations(files_to_merge, result_filename)

print("Merged!")

print("Removing free version Spire library 'watermarks/ads'...")
p=P(result_filename)
slides = p.slides
for slide in slides:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        tf = shape.text_frame
        if "Spire.Presentation for Python" in tf.text:
# remove annoying warning for not paying for spire lib
            sp = shape._element
            sp.getparent().remove(sp)
p.save(result_filename)

print("Done!")
print(f"Results saved to file: {ppt_to_merge_folder}/{result_filename}")
